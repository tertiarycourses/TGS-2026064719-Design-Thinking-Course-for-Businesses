#!/usr/bin/env python3
"""Generate the Design Thinking Course for Businesses Lesson Plan (LP) DOCX.

Cover page + Document Version Control Record + auto TOC + Arial 11pt body +
colour-coded 1-day schedule table (9:30am-6:30pm, 8 training hours, 1h lunch,
tea breaks within training time, final assessment from 5:30pm). Topics/activities
come from course_data + the domain data files so the LP stays aligned with the
deck, guide and labs.
"""
import os, sys
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
import course_data as C
from data_domain1 import DOMAIN1; from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3; from data_domain4 import DOMAIN4
ACT=DOMAIN1+DOMAIN2+DOMAIN3+DOMAIN4
import prodoc
def _find_repo(start):
    env=os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env): return env
    d=start
    for _ in range(8):
        d=os.path.dirname(d)
        if os.path.isdir(os.path.join(d,"courseware")) and os.path.isdir(os.path.join(d,"labs")): return d
    return os.path.dirname(os.path.dirname(HERE))
REPO=_find_repo(HERE); ASSETS=os.path.join(os.path.dirname(HERE),"assets")

BRAND=RGBColor(0x1F,0x6F,0xEB); DARK=RGBColor(0x11,0x18,0x27); GREY=RGBColor(0x55,0x5B,0x66)
HEADER_FILL="1F6FEB"; TOPIC_FILL="E8F0FE"; BREAK_FILL="FFF4E5"; LUNCH_FILL="FDE9D9"; ASSESS_FILL="E8F7EE"

def act_titles(nums):
    return "; ".join(f"Activity {a['num']}: {a['title']}" for a in ACT if a['num'] in nums)

# ------------------------------------------------ slide references
# Derived from the ACTUAL built deck so the LP can never disagree with the PPT.
# Slide numbers are deck positions (1-based), i.e. what PowerPoint shows in the
# slide sorter. Falls back to an empty map if the deck has not been built yet.
def _slide_map():
    """Scan the built PPTX and return {topic_num: 'Slides a–b'} by locating each
    'TOPIC nn' section divider and running to the slide before the next divider."""
    import glob
    try:
        from pptx import Presentation
    except ImportError:
        return {}
    hits=sorted(glob.glob(os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-*.pptx")))
    if not hits: return {}
    prs=Presentation(hits[-1])
    texts=[]
    for s in prs.slides:
        texts.append(" \n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame))
    starts={}
    for tp in C.TOPICS:
        marker=f"TOPIC {tp['code']}"
        for i,t in enumerate(texts,1):
            # the section divider is the slide whose FIRST line is exactly the marker
            first=t.strip().split("\n")[0].strip()
            if first==marker:
                starts[tp["num"]]=i; break
    # the closing section divider bounds the final topic
    close=len(texts)+1
    for i,t in enumerate(texts,1):
        if t.strip().split("\n")[0].strip()=="WRAP-UP":
            close=i; break
    out={}
    ordered=sorted(starts.items())
    for idx,(num,start) in enumerate(ordered):
        end=(ordered[idx+1][1]-1) if idx+1<len(ordered) else close-1
        out[num]=f"Slides {start}–{end}"
    return out
SLIDE_MAP=_slide_map()

# ------------------------------------------------ schedule (single source of truth for timing)
# (minutes, kind, activity_text)  kind: admin/topic/lab/break/lunch/assess/recap
# Start/end clock times are COMPUTED from the durations below so the table can
# never disagree with itself. Total course hours (everything except lunch) = 480.
DAY_START = 9*60 + 30   # 9:30 am

def _clock(m):
    return f"{(m//60)%24 if (m//60)%24<=12 else (m//60)%24-12}:{m%60:02d}"

def _expand(rows):
    """(mins, kind, text) → (start, end, mins, kind, text) with computed clock times."""
    out=[]; t=DAY_START
    for mins,kind,text in rows:
        out.append((_clock(t),_clock(t+mins),mins,kind,text)); t+=mins
    return out

SCHEDULE_RAW = {
 1: (C.DAY_THEMES[1], [
    (20,"admin","Welcome, trainer and learner introductions, ground rules, course outline, learning outcomes and mandatory SSG digital attendance (AM). Download of course material from the LMS/TMS portal and introduction to the browser-based ed-tools used in the activities"),
    (45,"topic","Topic 1 — Key Concepts & Principles of Design Thinking: what design thinking is, why it matters, ROI and business value, analytical vs creative thinking, the design thinker's mindset, growth mindset, and the desirability/viability/feasibility lenses (concepts + facilitated discussion)"),
    (15,"break","Tea break (within training time)"),
    (40,"lab","Hands-on: "+act_titles([1,2])),
    (55,"lab","Topic 2 — Applications of Design Thinking: corporate, service and public-sector use cases (Airbnb, IBM, Bank of America, Uber Eats, healthcare, Clean Team, GGRC, The Good Kitchen) and how to uncover opportunities. Hands-on: "+act_titles([3])),
    (60,"lunch","Lunch break"),
    (40,"topic","Topic 3 — Action Phases of Design Thinking: the Stanford d.school five-stage model — Empathize, Define, Ideate, Prototype and Test — including objectives, techniques and the iterative nature of the loop. Mandatory SSG digital attendance (PM)"),
    (40,"lab","Hands-on: "+act_titles([4])),
    (15,"break","Tea break (within training time)"),
    (35,"topic","Topic 4 — Methodologies and Visual Tools: empathy maps, personas and journey maps; POV and How Might We; brainstorming and idea-selection methods; low- and high-fidelity prototyping; user testing; and metrics for measuring design outcomes"),
    (30,"lab","Hands-on: "+act_titles([5])),
    (40,"lab","Hands-on: "+act_titles([6,7])),
    (10,"recap","Course recap, summary and Q&A, course feedback and the mandatory TRAQOM survey"),
    (5,"assess","Briefing for Assessment and Assessment digital attendance"),
    (70,"assess","Practical Performance (PP) — design-thinking tasks using the course ed-tools, 70 minutes, open book"),
    (20,"assess","Oral Questioning (OQ) — the assessor questions each learner verbally on the underpinning knowledge, 20 minutes, open book"),
 ]),
}
SCHEDULE = {d:(theme,_expand(rows)) for d,(theme,rows) in SCHEDULE_RAW.items()}

def _assert_timetable_matches_deck():
    """The deck's Lesson Plan slide prints C.DAY_START / LUNCH_START / LUNCH_END /
    DAY_END. Those are the same window this schedule computes, so verify they agree —
    otherwise learners are shown one timetable and taught to another."""
    def ampm(m):
        h=(m//60)%24; suffix="am" if h<12 else "pm"; hh=h if h<=12 else h-12
        return f"{hh}:{m%60:02d}{suffix}"
    t=DAY_START; found=None
    for mins,kind,_ in SCHEDULE_RAW[1][1]:
        if kind=="lunch": found=(ampm(t),ampm(t+mins))
        t+=mins
    end=ampm(t)
    want=(C.LUNCH_START,C.LUNCH_END)
    assert found==want, f"deck lunch {want} != schedule lunch {found} — fix course_data.LUNCH_*"
    assert end==C.DAY_END, f"deck day end {C.DAY_END} != schedule end {end} — fix course_data.DAY_END"
    assert ampm(DAY_START)==C.DAY_START, f"deck start {C.DAY_START} != schedule start {ampm(DAY_START)}"
_assert_timetable_matches_deck()

# ------------------------------------------------ build document
doc=Document()
normal=doc.styles["Normal"]; normal.font.name="Arial"; normal.font.size=Pt(11)
prodoc.style_headings(doc)

prodoc.add_cover_page(doc,"LESSON PLAN",C.TITLE,C.VERSION.lstrip("v"),
                      org_logo=os.path.join(ASSETS,"tertiary-infotech-logo.png"),
                      course_logo=None, course_code=C.COURSE_CODE)
prodoc.add_version_control(doc,[
 ("14.0","1 January 2024","Previous release of the 1-day lesson plan for the Design Thinking Course for Businesses, issued under the superseded course reference.",C.TRAINER),
 (C.VERSION.lstrip("v"),C.VERSION_DATE,
  f"Re-issued under WSQ course code {C.COURSE_CODE}. Content carried over from the v14 master trainer deck and restructured to the current Tertiary Infotech WSQ house standard. Added seven structured hands-on activities built on the browser-based ed-tools (Design Thinking Toolkit, RACI, Scrum, Digital Transformation and BCM), and aligned the schedule, slide references and assessment instruments (PP 70 minutes + OQ 20 minutes).",C.TRAINER),
])
prodoc.add_toc(doc)

def H(text,level=1):
    return doc.add_heading(text,level=level)

H("Course Information",1)
info=[("Course Title",C.TITLE),("WSQ Course Reference",C.COURSE_CODE),
      ("TSC Title / Code",f"{C.TSC_TITLE}  ·  {C.TSC_CODE}"),
      ("Training Provider",C.ORG+"  ("+C.UEN.replace('UEN: ','UEN ')+")"),
      ("Duration","1 day · 8 course hours (385 minutes instruction + 95 minutes assessment; the 1-hour lunch break is excluded)"),
      ("Daily Timing","9:30 am – 6:30 pm, including a 1-hour lunch break; tea breaks are taken within training time"),
      ("Mode","Instructor-led, facilitated hands-on design-thinking activities using browser-based ed-tools"),
      ("Assessment","Practical Performance (PP, 70 minutes) and Oral Questioning (OQ, 20 minutes), both open book"),
      ("Trainer",C.TRAINER)]
t=doc.add_table(rows=0,cols=2); t.style="Table Grid"
for k,v in info:
    c=t.add_row().cells; c[0].text=""; r=c[0].paragraphs[0].add_run(k); r.bold=True; r.font.size=Pt(10)
    prodoc._shade_cell(c[0],TOPIC_FILL)
    c[1].text=""; c[1].paragraphs[0].add_run(v).font.size=Pt(10)

H("Learning Outcomes",1)
doc.add_paragraph("On completion of this course, learners will be able to:")
for lo in C.LEARNING_OUTCOMES:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(lo).font.size=Pt(10.5)

H("Skills Framework Alignment",1)
doc.add_paragraph(f"This course is aligned to the Skills Framework for Design, TSC \"{C.TSC_TITLE}\" ({C.TSC_CODE}).")
p=doc.add_paragraph(); r=p.add_run("Abilities"); r.bold=True
for a in C.TSC_ABILITIES:
    pp=doc.add_paragraph(style="List Bullet"); pp.add_run(a).font.size=Pt(10.5)
p=doc.add_paragraph(); r=p.add_run("Knowledge"); r.bold=True
for k in C.TSC_KNOWLEDGE:
    pp=doc.add_paragraph(style="List Bullet"); pp.add_run(k).font.size=Pt(10.5)

H("Assessment",1)
for a in [C.ASSESSMENT["written"],C.ASSESSMENT["practical"],
          "Format: Open Book — course slides, Learner Guide and approved materials only.",
          "The final assessment is conducted at the end of Day 1, after the course recap and TRAQOM survey.",
          "Learners must be assessed as 'Competent' in both instruments to be awarded the Statement of Attainment.",
          C.ASSESSMENT["note"]]:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(a).font.size=Pt(10.5)

H("Ed-Tools Used in the Activities",1)
doc.add_paragraph("All activity tools are browser-based and require no installation:")
et=doc.add_table(rows=0,cols=3); et.style="Table Grid"
def set_cell(cell,text,bold=False,size=9.5,color=None,fill=None,align=None):
    cell.text=""; p=cell.paragraphs[0]
    if align: p.alignment=align
    r=p.add_run(text); r.bold=bold; r.font.size=Pt(size); r.font.name="Arial"
    if color: r.font.color.rgb=color
    if fill: prodoc._shade_cell(cell,fill)
hdr=et.add_row().cells
for i,h in enumerate(["Ed-Tool","Link","Used for"]):
    set_cell(hdr[i],h,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
for nm,url,desc in C.EDTOOLS:
    c=et.add_row().cells
    set_cell(c[0],nm,bold=True,size=9.5,fill=TOPIC_FILL); set_cell(c[1],url,size=9); set_cell(c[2],desc,size=9)
for row in et.rows:
    row.cells[0].width=Inches(1.7); row.cells[1].width=Inches(2.4); row.cells[2].width=Inches(2.7)

KIND_FILL={"topic":TOPIC_FILL,"break":BREAK_FILL,"lunch":LUNCH_FILL,"assess":ASSESS_FILL,
           "admin":"F3F5F8","recap":"F3F5F8","lab":None}

H("Course Schedule",1)
for day,(theme,rows) in SCHEDULE.items():
    H(f"Day {day} — {theme}",2)
    tbl=doc.add_table(rows=0,cols=3); tbl.style="Table Grid"; tbl.alignment=WD_TABLE_ALIGNMENT.CENTER
    hdr=tbl.add_row().cells
    for i,htext in enumerate(["Time","Duration","Topic / Activity"]):
        set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
    training=0; assess=0
    for start,end,mins,kind,text in rows:
        cells=tbl.add_row().cells; fill=KIND_FILL.get(kind)
        set_cell(cells[0],f"{start}–{end}",bold=(kind in ("topic","assess")),size=9.5,fill=fill)
        set_cell(cells[1],f"{mins} min",size=9.5,fill=fill)
        set_cell(cells[2],text,bold=(kind in ("topic","assess")),size=9.5,fill=fill)
        if kind=="assess": assess+=mins
        elif kind!="lunch": training+=mins
    for row in tbl.rows:
        row.cells[0].width=Inches(1.15); row.cells[1].width=Inches(0.9); row.cells[2].width=Inches(4.75)
    total=training+assess
    p=doc.add_paragraph()
    r=p.add_run(f"Total course hours: {total} minutes ({total//60} hours) — {training} minutes of instruction "
                f"(including tea breaks taken within training time) plus {assess} minutes of assessment. "
                f"The 1-hour lunch break is excluded.")
    r.italic=True; r.font.size=Pt(9.5); r.font.color.rgb=GREY
    assert total==480, f"Day {day} total = {total} minutes (instruction {training} + assessment {assess}), expected 480"

H("Activity Reference (aligned to topics and slides)",1)
tt=doc.add_table(rows=0,cols=4); tt.style="Table Grid"
hdr=tt.add_row().cells
for i,htext in enumerate(["Topic / Skills mapped","Slide positions","Activities","Learning outcome"]):
    set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
for tp in C.TOPICS:
    acts=[a for a in ACT if a["topic"]==tp["num"]]
    cells=tt.add_row().cells
    set_cell(cells[0],f"Topic {tp['code']}: {tp['title']}\n({tp['weighting']})",bold=True,size=9,fill=TOPIC_FILL)
    set_cell(cells[1],SLIDE_MAP.get(tp["num"],"—"),size=9,fill=TOPIC_FILL)
    set_cell(cells[2],", ".join(f"Activity {a['num']}" for a in acts) or "—",size=9)
    los=sorted({lo.split("/")[0].strip() for a in acts for lo in [a["objective"]]})
    set_cell(cells[3],"; ".join(los),size=9)
for row in tt.rows:
    row.cells[0].width=Inches(2.5); row.cells[1].width=Inches(1.0)
    row.cells[2].width=Inches(1.6); row.cells[3].width=Inches(1.7)

H("Detailed Activity Plan",1)
at=doc.add_table(rows=0,cols=4); at.style="Table Grid"
hdr=at.add_row().cells
for i,htext in enumerate(["Activity","Learners will…","Ed-tool","Trainer assesses"]):
    set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
for a in ACT:
    c=at.add_row().cells
    set_cell(c[0],f"Activity {a['num']}\n{a['title']}",bold=True,size=9,fill=TOPIC_FILL)
    set_cell(c[1],a["desc"],size=9)
    set_cell(c[2],a.get("edtool",""),size=8.5)
    set_cell(c[3],a["test"],size=9)
for row in at.rows:
    row.cells[0].width=Inches(1.5); row.cells[1].width=Inches(2.4)
    row.cells[2].width=Inches(1.3); row.cells[3].width=Inches(1.6)

H("Trainer Notes",1)
for n in [
 "The deck is deliberately visual and carries no step-by-step instruction slides — walk learners through the steps using the Learner Guide and the labs/ files while they work.",
 "Activity 1 (Design a Vase) is the key hook of the morning: run round 1 strictly to 2 minutes so the narrow-brief effect is felt before you reveal the reframed brief.",
 "Activity 4 (the Wallet Project) is the only full five-phase cycle. Keep each phase to time and enforce the swap so both partners design.",
 "Ed-tools are browser-based; verify learners can reach alfredang.github.io before the first activity and have a paper fallback ready.",
 "Collect evidence for the Practical Performance assessment as learners work — their empathy map, POV statement and prototype form the basis of the PP tasks.",
 "Digital attendance must be taken three times: AM, PM and Assessment. Remind learners of the 75% attendance rule for funding eligibility.",
]:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(n).font.size=Pt(10.5)

prodoc.add_page_numbers(doc)
prodoc.enable_update_fields(doc)
OUT=os.path.join(REPO,"courseware",f"LP-{C.SHORT_TITLE}.docx")
doc.save(OUT)
print("Saved",OUT)
