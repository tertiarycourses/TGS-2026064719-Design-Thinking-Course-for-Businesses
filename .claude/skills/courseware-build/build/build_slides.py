#!/usr/bin/env python3
"""Generate the Design Thinking Course for Businesses slide deck (all-white Tertiary house style).

HIGHLY VISUAL by design: tile grids, chevron flows, profile cards, stat bands,
quadrant matrices, canvas mock-ups and quote slides. Per the course brief the
deck deliberately carries NO step-by-step instruction slides — the detailed
steps live in the Learner Guide and the labs/ files. Each activity gets a
visual overview + a visual "what you'll produce" card instead.

Content is driven entirely by course_data.py + data_domainN.py so the deck stays
100% aligned with the LP, LG and labs.
"""
import os, sys, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4

def _find_repo(start):
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")

# ---------------- palette ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)
ROSE=RGBColor(0xE1,0x1D,0x48); MINT=RGBColor(0xE8,0xF7,0xEE)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def roundrect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(5,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.25)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=False
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    # auto-shrink so a long title never collides with the rule at 1.7"
    n=len(title)
    size=29 if n<=52 else (25 if n<=68 else (22 if n<=86 else 20))
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.78),[[(title,size,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge — five-phase design thinking mark
    bx=Inches(10.15); by=Inches(0.62)
    for i,(lab,col) in enumerate([("E",BLUE),("D",TEAL),("I",AMBER),("P",VIOLET),("T",ROSE)]):
        d=Inches(0.46)
        oval(s,int(bx+i*Inches(0.5)),by,d,d,col)
        txt(s,int(bx+i*Inches(0.5)),by,d,d,[[(lab,14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,bx,Inches(1.18),Inches(2.5),Inches(0.3),[[("DESIGN THINKING",9,GREY,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [(f"TSC: {C.TSC_TITLE}  ·  {C.TSC_CODE}",13,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",13,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(0.12),BLUE); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(0.12),TEAL)
    if lhead: txt(s,Inches(1.1),Inches(2.2),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.2),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.75),Inches(5.2),Inches(3.7),left,size=15)
    bullets(s,Inches(7.2),Inches(2.75),Inches(5.05),Inches(3.7),right,size=15,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.7),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(3.0),Inches(3.2),Inches(3.35),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
def quote_slide(quote,author,role="",color=VIOLET):
    """Large pull-quote — visual, no bullets."""
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.0),Inches(1.35),Inches(2),Inches(1.4),[[("“",130,RGBColor(0xE2,0xE8,0xF0),True)]])
    txt(s,Inches(1.55),Inches(2.15),Inches(10.6),Inches(2.9),[[(quote,30,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,Inches(1.6),Inches(5.25),Inches(1.8),Inches(0.06),color)
    runs=[[(author,18,color,True)]]
    if role: runs.append([(role,13,GREY,False)])
    txt(s,Inches(1.6),Inches(5.5),Inches(10),Inches(0.9),runs,space=4)
    footer(s); return s

PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,17,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE,labels=None):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.35); ch=Inches(3.5); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)] if labels else color
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),col)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,col)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        ty=int(y+Inches(1.45))
        if labels:
            txt(s,x+Inches(0.12),ty,cw-Inches(0.24),Inches(0.5),[[(labels[i],17,col,True)]],align=PP_ALIGN.CENTER)
            ty=int(y+Inches(2.0))
        txt(s,x+Inches(0.14),ty,cw-Inches(0.28),int(ch-(ty-y)-Inches(0.2)),[[(st,13,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,col,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def stat_band(title,stats,source,kicker=None,color=TEAL):
    """Big-number statistic band — highly visual proof points."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(stats); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.3)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.3); ch=Inches(3.3)
    for i,(num,cap) in enumerate(stats):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.12),col)
        txt(s,x,int(y+Inches(0.75)),cw,Inches(1.3),[[(num,54,col,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.2),int(y+Inches(2.1)),cw-Inches(0.4),Inches(1.0),[[(cap,15,INK,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(5.95),Inches(11.6),Inches(0.4),[[(source,11,GREY,False)]])
    footer(s); return s
def venn3(title,circles,centre,kicker=None):
    """Three overlapping circles — desirability / viability / feasibility.

    Geometry: circle centres sit on an equilateral triangle around (cx,cy);
    each oval is positioned by its own centre so the three overlap evenly and
    the whole figure stays inside the content area."""
    s=head(slide(),title,kicker,kcolor=VIOLET)
    cx=Inches(6.67); cy=Inches(4.25); R=Inches(1.42); d=Inches(0.92)  # d = centre offset
    cols=[BLUE,TEAL,AMBER]
    cent=[(int(cx-d),int(cy-d*0.55)),          # desirability — upper left
          (int(cx+d),int(cy-d*0.55)),          # viability    — upper right
          (int(cx),  int(cy+d*0.95))]          # feasibility  — lower centre
    for i,(ox,oy) in enumerate(cent):
        sp=oval(s,int(ox-R),int(oy-R),int(R*2),int(R*2),cols[i])
        sp.fill.background(); sp.line.color.rgb=cols[i]; sp.line.width=Pt(2.5)
    # centre label sits at the true centroid of the three circles
    gx=int(sum(p[0] for p in cent)/3); gy=int(sum(p[1] for p in cent)/3)
    txt(s,int(gx-Inches(1.0)),int(gy-Inches(0.2)),Inches(2.0),Inches(0.45),
        [[(centre,15,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # labelled cards down the left, colour-matched to each circle
    for i,(nm,desc) in enumerate(circles):
        y=int(Inches(2.05)+i*Inches(1.62))
        rect(s,Inches(0.85),y,Inches(3.2),Inches(1.3),LIGHT)
        rect(s,Inches(0.85),y,Inches(0.1),Inches(1.3),cols[i])
        txt(s,Inches(1.12),y,Inches(2.8),Inches(1.3),
            [[(nm,15,cols[i],True)],[(desc,11.5,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def quadrant(title,xlab,ylab,quads,kicker=None,color=BLUE):
    """2x2 matrix — e.g. business value vs novelty."""
    s=head(slide(),title,kicker,kcolor=color)
    X0=Inches(2.6); Y0=Inches(2.05); W=Inches(8.6); H=Inches(4.3)
    hw=int(W/2); hh=int(H/2)
    fills=[RGBColor(0xEE,0xF4,0xFF),RGBColor(0xE8,0xF7,0xEE),RGBColor(0xF6,0xF7,0xF9),RGBColor(0xFF,0xF6,0xE5)]
    for i,(qt,qd) in enumerate(quads):
        r=i//2; c=i%2
        x=int(X0+hw*c); y=int(Y0+hh*r)
        rect(s,x,y,hw,hh,fills[i],line=LINE)
        txt(s,x+Inches(0.22),y+Inches(0.2),hw-Inches(0.44),Inches(0.5),[[(qt,16,INK,True)]])
        txt(s,x+Inches(0.22),y+Inches(0.75),hw-Inches(0.44),Inches(1.3),[[(qd,12,GREY,False)]])
    txt(s,X0,int(Y0+H+Inches(0.12)),W,Inches(0.4),[[(xlab,13,color,True)]],align=PP_ALIGN.CENTER)
    tb=txt(s,Inches(0.85),int(Y0+H/2-Inches(0.25)),Inches(1.6),Inches(0.5),[[(ylab,13,color,True)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,13,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def download_slide():
    """Visual 'how to download your course material' — mock browser, not a bare link."""
    s=head(slide(),"Download Your Course Material","COURSE PORTAL  ·  LMS / TMS",kcolor=BLUE)
    bx=Inches(0.85); by=Inches(2.0); bw=Inches(7.2); bh=Inches(4.4)
    rect(s,bx,by,bw,bh,WHITE,line=LINE)
    rect(s,bx,by,bw,Inches(0.52),RGBColor(0xF1,0xF5,0xF9))
    for i,col in enumerate([ROSE,AMBER,TEAL]):
        oval(s,int(bx+Inches(0.22)+i*Inches(0.3)),int(by+Inches(0.17)),Inches(0.18),Inches(0.18),col)
    rect(s,int(bx+Inches(1.25)),int(by+Inches(0.12)),int(bw-Inches(1.55)),Inches(0.3),WHITE,line=LINE)
    txt(s,int(bx+Inches(1.42)),int(by+Inches(0.12)),int(bw-Inches(1.7)),Inches(0.3),
        [[("🔒  lms-tms.tertiaryinfotech.com",11,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    # portal rows
    rows=[("Trainer Slides","PPT",BLUE),("Learner Slides","PDF",TEAL),
          ("Learner Guide","PDF",VIOLET),("Lesson Plan","PDF",AMBER),("Activities / Labs","FOLDER",ROSE)]
    ry=int(by+Inches(0.95))
    txt(s,int(bx+Inches(0.35)),int(by+Inches(0.62)),Inches(6),Inches(0.35),[[("My Courses  ›  Design Thinking for Businesses",11,INK,True)]])
    for i,(nm,kind,col) in enumerate(rows):
        y=int(ry+i*Inches(0.63))
        rect(s,int(bx+Inches(0.35)),y,int(bw-Inches(0.7)),Inches(0.52),LIGHT)
        rect(s,int(bx+Inches(0.35)),y,Inches(0.08),Inches(0.52),col)
        txt(s,int(bx+Inches(0.6)),y,Inches(3.6),Inches(0.52),[[(nm,13,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
        rect(s,int(bx+bw-Inches(1.72)),int(y+Inches(0.11)),Inches(1.3),Inches(0.3),col)
        txt(s,int(bx+bw-Inches(1.72)),int(y+Inches(0.11)),Inches(1.3),Inches(0.3),
            [[("⤓  "+kind,10,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # side steps
    sx=Inches(8.45); sw=Inches(4.05)
    for i,(t1,t2) in enumerate([("Sign in","Use the e-mail you registered with"),
                                 ("Open My Courses","Select this course from your list"),
                                 ("Click ⤓ Download","Slides, Learner Guide and Lesson Plan"),
                                 ("Keep them open","Your assessment is open book")]):
        y=int(by+i*Inches(1.13))
        rect(s,sx,y,sw,Inches(0.98),LIGHT); rect(s,sx,y,Inches(0.09),Inches(0.98),PALETTE[i%4])
        oval(s,int(sx+Inches(0.28)),int(y+Inches(0.26)),Inches(0.46),Inches(0.46),PALETTE[i%4])
        txt(s,int(sx+Inches(0.28)),int(y+Inches(0.26)),Inches(0.46),Inches(0.46),
            [[(str(i+1),14,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,int(sx+Inches(0.92)),y,int(sw-Inches(1.1)),Inches(0.98),
            [[(t1,13,INK,True)],[(t2,10.5,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=2)
    footer(s); return s
# Post-course framing for each ed-tool — used on the closing slide so it is a genuine
# call-to-action rather than a repeat of the front-of-course orientation slide.
EDTOOL_TAKEAWAY=[
 "Re-run the Empathy Map and POV canvases on a live problem with your own team.",
 "Assign an owner for your first initiative so it survives the handover to delivery.",
 "Turn the concept you prototyped today into a backlog and a first sprint goal.",
 "Position the initiative inside your organisation's wider transformation roadmap.",
 "Stress-test the solution for operational risk before you commit to rollout.",
]
def edtool_slide(title="Your Lab Ed-Tools",kicker="BROWSER-BASED  ·  NO INSTALL REQUIRED",
                 descs=None):
    """Visual index of the browser-based ed-tools. `descs` overrides the tile captions so
    the closing slide can carry post-course actions instead of repeating the intro slide."""
    s=head(slide(),title,kicker,kcolor=TEAL)
    n=len(C.EDTOOLS); X0=Inches(0.85); TOTW=Inches(11.63)
    # size the card block to the space between the rule (1.7") and the footer
    # (7.05"), leaving 0.25" clearance so the last card can never touch it.
    AREA=Inches(6.80)-Inches(1.95)
    gy=Inches(0.13); ch=int((AREA-gy*(n-1))/n)
    for i,(nm,url,desc) in enumerate(C.EDTOOLS):
        y=int(Inches(1.95)+i*(ch+gy)); col=PALETTE[i%len(PALETTE)]
        rect(s,X0,y,TOTW,ch,LIGHT); rect(s,X0,y,Inches(0.1),ch,col)
        bd=Inches(0.48); by=int(y+ch/2-bd/2)
        oval(s,int(X0+Inches(0.3)),by,bd,bd,col)
        txt(s,int(X0+Inches(0.3)),by,bd,bd,
            [[(str(i+1),15,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,int(X0+Inches(0.98)),y,Inches(3.5),ch,[[(nm,14,INK,True)],[(url.replace("https://",""),9.5,col,False)]],
            anchor=MSO_ANCHOR.MIDDLE,space=2)
        cap=descs[i] if descs else desc
        txt(s,int(X0+Inches(4.6)),y,int(TOTW-Inches(4.85)),ch,[[(cap,12,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def canvas_slide(title,quads,kicker,centre=None,color=BLUE):
    """Mock-up of a 4-quadrant canvas (empathy map etc) — visual, not bullets."""
    s=head(slide(),title,kicker,kcolor=color)
    X0=Inches(2.0); Y0=Inches(2.0); W=Inches(9.3); H=Inches(4.4)
    hw=int(W/2); hh=int(H/2)
    cols=[BLUE,TEAL,VIOLET,AMBER]
    for i,(qt,qd) in enumerate(quads):
        r=i//2; c=i%2
        x=int(X0+hw*c); y=int(Y0+hh*r)
        rect(s,x,y,hw,hh,LIGHT,line=WHITE)
        rect(s,x,y,hw,Inches(0.1),cols[i])
        # bottom-row headings are nudged down so the centre badge never covers them
        hy=y+Inches(0.25) if r==0 else y+Inches(0.62)
        txt(s,x+Inches(0.28),hy,hw-Inches(0.56),Inches(0.5),[[(qt,17,cols[i],True)]])
        txt(s,x+Inches(0.28),hy+Inches(0.55),hw-Inches(0.56),Inches(1.3),[[(qd,12.5,GREY,False)]])
    if centre:
        d=Inches(1.05)
        sp=oval(s,int(X0+W/2-d/2),int(Y0+H/2-d/2),d,d,color)
        sp.line.color.rgb=WHITE; sp.line.width=Pt(3)
        txt(s,int(X0+W/2-d/2),int(Y0+H/2-d/2),d,d,[[(centre,12,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker,edtool=None):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,20,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(0.1),BLUE)
    txt(s,Inches(1.1),Inches(4.55),Inches(11),Inches(0.4),[[("You'll produce",13,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.92),Inches(11),Inches(0.6),[[(build,17,INK,True)]])
    txt(s,Inches(1.1),Inches(5.62),Inches(11.2),Inches(0.6),[[("Ed-tools:  ",12,GREY,True),(services,12,GREY,False)]])
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),MINT)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(0.1),TEAL)
    txt(s,Inches(1.2),Inches(2.65),Inches(11),Inches(0.5),[[("✅  You're done when…",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.35),Inches(11),Inches(1.4),[[(text,17,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
cover()

# ---------------- ADMIN (front) ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
def attendance_slide():
    """Digital attendance as a 3-step flow + a rule callout — not a bullet wall."""
    s=head(slide(),"Digital Attendance (Mandatory)","TRAQOM · SSG DIGITAL ATTENDANCE",kcolor=BLUE)
    steps=[("Trainer shows the QR","The trainer or administrator displays the digital attendance QR code generated from the SSG portal."),
           ("Scan with your phone","Open your mobile phone camera and scan the QR code shown on screen."),
           ("Submit — three times","Attendance is taken at AM, PM and again before the Assessment.")]
    X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.32); n=len(steps)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.1); ch=Inches(2.85); bd=Inches(0.85)
    for i,(t1,t2) in enumerate(steps):
        x=int(X0+(cw+gap)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),col)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.4)),bd,bd,col)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.4)),bd,bd,[[(str(i+1),30,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.2),int(y+Inches(1.45)),cw-Inches(0.4),Inches(0.45),[[(t1,16,col,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.2),int(y+Inches(1.95)),cw-Inches(0.4),Inches(0.8),[[(t2,12.5,GREY,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,col,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,X0,Inches(5.35),TOTW,Inches(1.15),MINT); rect(s,X0,Inches(5.35),Inches(0.1),Inches(1.15),TEAL)
    txt(s,int(X0+Inches(0.35)),Inches(5.35),TOTW-Inches(0.7),Inches(1.15),
        [[("⚠  75% minimum attendance",16,RGBColor(0x12,0x7A,0x3E),True)],
         [("Mandatory for all WSQ-funded courses. You must meet the 75% attendance rate based on the SSG Digital Attendance record to be eligible for assessment and funding.",12.5,INK,False)]],
        anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
attendance_slide()
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Qualifications","PhD; ACTA/ACLP-certified adult educator with design thinking and innovation facilitation experience."),
  ("Delivers","WSQ courses on design thinking, innovation, digital transformation and emerging technology."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name, organisation and role.",
 "A product or service experience you love — and one that frustrates you.",
 "What you want to be able to design or improve after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
download_slide()
edtool_slide()
tile_grid("Skills Framework Alignment",[
 ("TSC Title",C.TSC_TITLE),("TSC Code",C.TSC_CODE),
 ("Abilities assessed","A1–A6: apply methodologies, uncover opportunities, use metrics, embed practice, facilitate prototypes, communicate outcomes."),
 ("Knowledge assessed","K1–K7: concepts, importance, traits, use cases, organisational approaches, methods and metrics.")],
 kicker="SKILLS FRAMEWORK FOR DESIGN",cols=2,size=14,accent=VIOLET)
two_col("Lesson Plan — 1 Day, 8 Course Hours",[
 ("Morning — 9:30am to 12:45pm",0),
 ("Digital Attendance (AM), introductions, learning outcomes",1),
 ("Topic 1: Key Concepts & Principles (Activities 1–2)",1),
 ("Topic 2: Applications of Design Thinking (Activity 3)",1),
 ("Tea break taken within training time",1),
 ("Lunch break 12:45pm – 1:45pm",1)],
 [("Afternoon — 1:45pm to 6:30pm",0),
 ("Digital Attendance (PM)",1),
 ("Topic 3: Action Phases — the Wallet Project (Activity 4)",1),
 ("Topic 4: Methodologies & Visual Tools (Activities 5–7)",1),
 ("Tea break taken within training time",1),
 ("Course feedback, TRAQOM survey and Final Assessment",1)],
 kicker="SCHEDULE",lhead="Morning",rhead="Afternoon")
tile_grid("Learning Outcomes",[
 ("LO1 · Understand","Key concepts of design thinking, to communicate design outcomes."),
 ("LO2 · Apply","Design thinking to generate new ideas for the organisation."),
 ("LO3 · Uncover","Opportunities for applying design thinking."),
 ("LO4 · Implement","Plans to embed the stages of design thinking across the organisation."),
 ("LO5 · Execute","The design concept through prototypes."),
 ("LO6 · Measure","Outcomes of design ideas and solutions using metrics.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=14)
def briefing_slide():
    tile_grid("Briefing for Assessment",[
     ("Phones away","Place phones and other materials under the table or on the floor."),
     ("No photos","No photographs or recording of the assessment scripts."),
     ("No discussion","Work individually — no discussion during the assessment."),
     ("Black or blue pen","Use a black or blue pen for hard-copy assessments."),
     ("No correction fluid","No liquid paper or correction tape may be used."),
     ("Time is up","Scripts are collected when the time is up.")],
     kicker="BEFORE YOU SIT THE ASSESSMENT",cols=2,size=14,accent=AMBER)
def assessment_slide():
    """Two instrument cards + an open-book / competency callout."""
    s=head(slide(),"Assessment","FINAL ASSESSMENT",kcolor=BLUE)
    cards=[(BLUE,"Written Assessment (WA)","Short-Answer Questions (SAQ)",
            f"{C.WA_MINUTES} minutes  ·  Open book",
            "Tests your KNOWLEDGE of design thinking concepts, phases, tools and metrics."),
           (TEAL,"Practical Performance (PP)","Design-thinking tasks",
            f"{C.PP_MINUTES} minutes  ·  Open book",
            "Tests your ABILITY to apply the methods using the course ed-tools.")]
    X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.35)
    cw=int((TOTW-gap)/2); y=Inches(2.0); ch=Inches(3.15)
    for i,(col,t1,t2,t3,t4) in enumerate(cards):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.12),col)
        txt(s,x+Inches(0.35),int(y+Inches(0.35)),cw-Inches(0.7),Inches(0.55),[[(t1,20,col,True)]])
        txt(s,x+Inches(0.35),int(y+Inches(1.0)),cw-Inches(0.7),Inches(0.4),[[(t2,15,INK,True)]])
        rect(s,int(x+Inches(0.35)),int(y+Inches(1.55)),Inches(2.9),Inches(0.38),col)
        txt(s,int(x+Inches(0.35)),int(y+Inches(1.55)),Inches(2.9),Inches(0.38),
            [[(t3,12,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.35),int(y+Inches(2.15)),cw-Inches(0.7),Inches(0.85),[[(t4,13,GREY,False)]])
    rect(s,X0,Inches(5.45),TOTW,Inches(1.1),MINT); rect(s,X0,Inches(5.45),Inches(0.1),Inches(1.1),TEAL)
    txt(s,int(X0+Inches(0.35)),Inches(5.45),TOTW-Inches(0.7),Inches(1.1),
        [[("📖  Open book  ·  you must be assessed 'Competent' in both instruments",15,RGBColor(0x12,0x7A,0x3E),True)],
         [("Permitted materials: the course slides, the Learner Guide and approved materials only. A minimum of 75% attendance is required to be eligible. An appeal process is available if required.",12,INK,False)]],
        anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
briefing_slide()
assessment_slide()
flow_h("Assessment Flow",[
 "Scan the QR code on the LMS",
 "Scan the SSG QR code",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY",
 labels=["TRAQOM","Attendance","Assessment","Submit","Sign off"])
content("Criteria for Funding",[
 "A minimum attendance rate of 75%, based on the SSG Digital Attendance record.",
 "Complete the assessment and be assessed as 'Competent'.",
 "Complete the mandatory TRAQOM survey on the LMS/TMS portal."],kicker="ELIGIBILITY")

# ---------------- TOPICS ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]

def topic_extras(num):
    """Extra visual concept slides per topic — carried from the v14 master deck."""
    if num==1:
        quote_slide("You've got to start with the customer experience and work backwards to the technology.",
                    "Steve Jobs","Co-founder, Apple Inc.",color=BLUE)
        tile_grid("What is Design Thinking?",[
         ("A problem-solving approach","Aims to improve people's experience, not just ship a product."),
         ("Human-centred","Starts from a deep understanding of customers' needs and wants."),
         ("Creative and wide","Encourages consideration of a wide array of innovative solutions."),
         ("A mindset, not only a process","How you think about the problem matters as much as the steps you follow.")],
         kicker="TOPIC 01 · CONCEPT",cols=2,size=15)
        quote_slide("Design thinking is a human-centered approach to innovation that integrates the needs of people, the possibilities of technology, and the requirements for business success.",
                    "Tim Brown","British industrial designer · President & CEO of IDEO",color=TEAL)
        tile_grid("Benefits of Design Thinking",[
         "Offers a problem-solving roadmap","Helps unlock innovative ideas",
         "Solves real human problems","Creates behavioural change",
         "Increases customer satisfaction","Gains competitive advantage"],
         kicker="WHY IT MATTERS",cols=2,size=16,accent=TEAL)
        stat_band("The ROI of Design Thinking",C.ROI_STATS,C.ROI_SOURCE,
                  kicker="DESIGN-LED ORGANISATIONS OUTPERFORM",color=TEAL)
        two_col("Analytical and Creative Thinking",
         [("Left brain",0),("Analytical",1),("Rational",1),("Objective",1),
          ("Present & past",1),("Facts",1),("Order and pattern",1),("Planned",1)],
         [("Right brain",0),("Creative",1),("Holistic",1),("Subjective",1),
          ("Present & future",1),("Feelings",1),("Spatial",1),("Spontaneous",1)],
         kicker="DESIGN THINKING USES BOTH",lhead="Analysis",rhead="Creativity")
        two_col("Mindset — Traditional vs Design Thinker",
         [("Traditional thinker",0),
          ("\"We have this problem — let's get in a room and brainstorm solutions.\"",1),
          ("\"Our competitors just launched X; how can we do X quickly?\"",1),
          ("\"We have this technology — what can we use it for?\"",1),
          ("Starts from the business or the technology",1)],
         [("Design thinker",0),("Think users first",1),("Ask the right questions",1),
          ("Believe you can sketch",1),("Commit to explore",1),
          ("Prototype to test and evaluate",1),("Starts from the human being served",1)],
         kicker="THE SHIFT",lhead="Business-centric",rhead="Customer-centric")
        big_statement("Design thinking is a growth mindset.",
                      "Talent and ability are a starting point, not a ceiling — the same belief that lets a team keep iterating instead of defending its first idea.",
                      "MINDSET",color=VIOLET)
        venn3("Design Thinking for the Organisation",
              [("DESIRABILITY","Do our customers want it?"),
               ("VIABILITY","Does it work for the business?"),
               ("FEASIBILITY","Can technology deliver it?")],
              "INNOVATION",kicker="THE THREE LENSES")
    elif num==2:
        cards3("Design Thinking in Business",[
         (BLUE,"Airbnb",["Culture of experimentation","Design-led iteration","Failing start-up → billion-dollar business"]),
         (TEAL,"IBM",["IBM Design Thinking at scale","Complex teams and organisations","Restructured how products are built"]),
         (VIOLET,"Bank of America",["Partnered with IDEO","Researched real savers","'Keep the Change' account"])],
         kicker="CORPORATE USE CASES")
        cards3("Design Thinking in Services & Society",[
         (BLUE,"Uber Eats",["Designs around observed journeys","Customer, courier and restaurant","Continuous field research"]),
         (TEAL,"Healthcare",["Stanford Hasso Plattner Institute","Redesigned the ER experience","Reduced patient anxiety"]),
         (VIOLET,"Public & social",["Clean Team — in-home toilets, Ghana","Golden Gate Regional Center","The Good Kitchen, Denmark"])],
         kicker="SERVICE & PUBLIC SECTOR USE CASES")
        tile_grid("Uncovering Opportunities in Your Organisation",[
         ("Look for friction","Where do customers or staff visibly struggle, complain or work around the system?"),
         ("Follow the complaints","Repeated complaints and support tickets are free user research."),
         ("Watch the workarounds","A workaround is a user telling you the design is wrong."),
         ("Find the drop-off","Where do people abandon the journey — and what happened just before?"),
         ("Ask the frontline","The people serving customers already know where the pain is."),
         ("Prioritise","Rank opportunities by impact and effort before committing a team.")],
         kicker="WHERE TO START",cols=2,size=13,accent=AMBER)
    elif num==3:
        flow_h("The Five Action Phases",
         ["Understand the experience of the user you are designing for.",
          "Synthesise findings into an actionable problem statement.",
          "Generate a wide variety of possible solutions.",
          "Build a tangible representation of your idea.",
          "Put it in users' hands and gather feedback."],
         kicker="STANFORD d.school MODEL",labels=["Empathize","Define","Ideate","Prototype","Test"])
        tile_grid("Action Phase 1 — Empathize",[
         ("Objective","Understand the experience, situation and emotion of the user you are designing for."),
         ("Observe","View users and their behaviour in the context of their lives. Don't judge."),
         ("Engage","Interact with people in conversations and interviews. Ask why — repeatedly."),
         ("Immerse","Experience what your user experiences, first-hand."),
         ("Beginner's mindset","Leave your assumptions behind; question what you think you already know."),
         ("Example — the MRI scanner","Immersion revealed that an MRI room terrifies a child. The redesign turned the scan into an adventure.")],
         kicker="TOPIC 03 · PHASE 1",cols=2,size=13,accent=BLUE)
        tile_grid("Action Phase 2 — Define",[
         ("Objective","Process and synthesise the findings into a problem statement you will address."),
         ("User","Develop an understanding of the persona you are designing for."),
         ("Needs","Select a limited set of needs to fulfil — needs should be verbs."),
         ("Insights","Express the insights you developed and define your design principles."),
         ("A good problem statement","Is human-centred, broad enough for creativity, narrow enough to be manageable."),
         ("Methods","Clustering, empathy mapping, POV statements, 'How Might We', why-how laddering.")],
         kicker="TOPIC 03 · PHASE 2",cols=2,size=13,accent=TEAL)
        quote_slide("A problem well-stated is a problem half-solved.",
                    "Charles F. Kettering","Inventor · Head of research, General Motors",color=TEAL)
        tile_grid("Action Phase 3 — Ideate",[
         ("Objective","Focus on idea generation — translate problems into solutions."),
         ("Go wide","Explore a large quantity and variety of ideas to get beyond the obvious."),
         ("Creativity","Combine rational thought with imagination; defer judgement."),
         ("Group synergy","Leverage the group to reach new ideas and build on others'."),
         ("Diverge, then converge","Separate the generation of ideas from their evaluation."),
         ("Facilitate actively","Create a curious, courageous and concentrated atmosphere.")],
         kicker="TOPIC 03 · PHASE 3",cols=2,size=13,accent=VIOLET)
        quote_slide("If at first the idea is not absurd, then there is no hope for it.",
                    "Albert Einstein","Theoretical physicist",color=VIOLET)
        tile_grid("Action Phase 4 — Prototype",[
         ("Objective","Build to think — a simple, cheap and fast way to shape ideas you can interact with."),
         ("It can be anything","A wall of post-its, a role-play, a space, an object, an interface, a storyboard."),
         ("Low fidelity","Basic models using simple materials — fast, cheap, easy to change."),
         ("High fidelity","Closer to the finished product — used later, when the concept is settled."),
         ("Just start building","Don't spend too much time; remember what you're testing for."),
         ("Build with the user in mind","The prototype exists to provoke a reaction, not to impress.")],
         kicker="TOPIC 03 · PHASE 4",cols=2,size=13,accent=AMBER)
        tile_grid("Action Phase 5 — Test",[
         ("Objective","Ask for feedback, learn about your user, reframe your POV and refine the prototype."),
         ("Show, don't tell","Put it in their hands and let them use it. Avoid over-explaining."),
         ("Create experiences","Let people talk about how they experience it and how they feel."),
         ("Compare","Let users test multiple prototypes to reveal latent needs and preferences."),
         ("Test in context","Use the natural setting where the product would really be used."),
         ("Negative feedback is data","Difficulty reveals new problems — revisit your solutions and reframe.")],
         kicker="TOPIC 03 · PHASE 5",cols=2,size=13,accent=ROSE)
        quote_slide("Insanity is doing the same thing over and over again and expecting different results.",
                    "Albert Einstein","On why testing must actually change the design",color=ROSE)
        big_statement("The five phases are a loop, not a line.",
                      "Insights from Test routinely send you back to Empathize or Define. Iteration is the method, not a sign of failure.",
                      "ITERATE",color=BLUE)
    elif num==4:
        tile_grid("Methods by Action Phase",[
         ("Empathize","Empathy map · Persona map · Journey map · Interviews · Bodystorming"),
         ("Define","POV statement · How Might We · Why-how laddering · Business Model Canvas"),
         ("Ideate","Brainstorm · Mind map · Sketchstorm · Analogies · Gamestorming · Dot voting"),
         ("Prototype","Sketches · Storyboards · Wireframes · Paper and 3D models"),
         ("Test","User test scripts · Feedback grids · Evaluation matrix · A/B comparison"),
         ("Measure","KPIs · CSAT and NPS · Activity metrics · Business value vs novelty")],
         kicker="TOPIC 04 · THE TOOLKIT",cols=2,size=13,accent=VIOLET)
        canvas_slide("The Empathy Map",
         [("SAYS","Direct, verbatim quotes from the interview. Don't paraphrase — the exact words carry the insight."),
          ("THINKS","What the user is thinking but not saying out loud. Worries, doubts, judgements."),
          ("DOES","The observable actions: what they actually do, where, and how long it takes."),
          ("FEELS","The emotional state and its intensity: anxious, rushed, proud, relieved.")],
         kicker="EMPATHIZE TOOL",centre="USER",color=BLUE)
        tile_grid("Reading an Empathy Map",[
         ("Look for contradictions","A positive quote beside a negative feeling is where the real insight hides."),
         ("Users are complex","Juxtaposition between quadrants is normal — and extremely useful."),
         ("Quote, don't summarise","Verbatim language keeps the user's meaning intact."),
         ("Turn it into a persona","Distil the map into a named persona with goals, frustrations and a defining quote.")],
         kicker="GUIDELINES",cols=2,size=14,accent=BLUE)
        tile_grid("Persona & Journey Mapping",[
         ("Persona mapping","Identify who your customers are and how they make decisions, then design for that specific person."),
         ("Journey mapping","Chart every touchpoint over time — before, during and after the core interaction."),
         ("Find the pain points","The journey map exposes where the experience breaks down."),
         ("Find the moments that matter","Not every touchpoint is equal — some disproportionately shape the whole experience.")],
         kicker="EMPATHIZE TOOLS",cols=2,size=14,accent=BLUE)
        flow_h("Point of View — The POV Formula",
         ["The specific person you are designing for — not 'everyone'.",
          "What they need to be able to do. Always a verb, never a solution.",
          "The surprising discovery that explains the need."],
         kicker="DEFINE TOOL",labels=["USER","NEED","INSIGHT"],color=TEAL)
        tile_grid("POV Examples",[
         ("Example 1","A busy working mom needs a way to bring bags with her when away from home, because she buys the equivalent of three dog-poop bags a day."),
         ("Example 2","A teacher needs a way to clean and store used sandwich bags at school, because he discards a new bag each morning rather than washing it."),
         ("The template","[USER] needs a way to [USER'S NEED] because [INSIGHT]."),
         ("State your assumptions","Write down what must be true for the POV to hold — then test it.")],
         kicker="DEFINE TOOL",cols=2,size=13,accent=TEAL)
        tile_grid("Ideation Methods",[
         ("Generate","Brainstorm · Mind map · Sketchstorm · Storyboard · Analogies · Provocation"),
         ("Go further","Bodystorm · Gamestorming · Cheatstorm · Crowdstorm · Co-creation workshops"),
         ("Brainstorm rules","Set a mission · Limit the time · Defer judgement · Build on ideas · Go for volume"),
         ("Select the best","Dot voting · Four categories · Bingo selection · Affinity maps · Now-Wow-How matrix"),
         ("Remember","10 ideas beat 3. 200 ideas beat 50. Quantity breeds quality."),
         ("Mind mapping","Start from the central problem and branch outwards to see connections you'd otherwise miss.")],
         kicker="IDEATE TOOLS",cols=2,size=13,accent=VIOLET)
        two_col("Low- vs High-Fidelity Prototyping",
         [("Low fidelity",0),("Paper sketches and storyboards",1),("Cardboard, foam, clay, building blocks",1),
          ("Fast, cheap, easy to throw away",1),("Use early, to test the concept",1),
          ("Invites honest criticism — it clearly isn't finished",1)],
         [("High fidelity",0),("Wireframes and interactive mock-ups",1),("Near-final materials and finishes",1),
          ("Slower and more expensive",1),("Use later, to test the detail",1),
          ("Risk: people critique the polish, not the idea",1)],
         kicker="PROTOTYPE TOOLS",lhead="Start here",rhead="Later")
        tile_grid("Testing Your Prototype",[
         ("Prototype as if you're right…","…but test as if you know you're wrong."),
         ("Find the cheapest test","What is the fastest, cheapest test that could disprove your key assumption?"),
         ("Define success up front","What metric will tell you it worked — pre-orders, votes, completion rate, smiles?"),
         ("Let users compare","Offer alternatives so the user can express a preference, not just be polite."),
         ("Ask them to think aloud","Narration reveals confusion that observation alone will miss."),
         ("Then revise","Revise and advance the idea until it genuinely meets the need.")],
         kicker="TEST TOOLS",cols=2,size=13,accent=ROSE)
        tile_grid("Measuring the Outcome of Design Ideas",[
         ("Traditional KPIs","Increased sales, ROI per project and other financial measures."),
         ("Customer feedback","Customer satisfaction, Net Promoter Score, campaign response, usability metrics."),
         ("Design thinking activity","Number of projects run, people trained, coaches developed."),
         ("Quick results","Concepts finished, projects launched, funded or in development."),
         ("Beyond execution","Track creative behaviours, not only delivery throughput."),
         ("Three business drivers","Understand customers better · protect share from disruption · develop innovative team dynamics.")],
         kicker="MEASURE OUTCOMES  ·  A3",cols=2,size=13,accent=AMBER)
        quadrant("Measuring Business Value and Novelty","NOVELTY  →","VALUE  →",
         [("Valuable · Novel","The target quadrant — genuinely new and worth doing. Protect and resource these."),
          ("Valuable · Not novel","Solid incremental improvement. Reliable returns, little differentiation."),
          ("Not valuable · Novel","Interesting but unprofitable. A nice demo, not a business case."),
          ("Not valuable · Not novel","Stop. This consumes capacity without creating value.")],
         kicker="PRIORITISING DESIGN IDEAS",color=AMBER)

for t in C.TOPICS:
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=f"SKILLS MAPPED: {t['weighting']}", cols=2, size=13)
    topic_extras(t["num"])
    acts=TOPIC_ACTS[t["num"]]
    # per activity — VISUAL overview + outcome only (no step-by-step slides by design)
    for a in acts:
        activity_overview(f"ACTIVITY {a['num']}", a["title"], a["desc"], a["build"], a["services"],
                          kicker=f"TOPIC {t['code']} · HANDS-ON",edtool=a.get("edtool"))
        test_slide(a["title"], a["test"], kicker=f"ACTIVITY {a['num']} · CHECK YOUR WORK")
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"].split("—",1)[-1].strip() for a in acts],
            kicker="TOPIC RECAP", size=17)
    if t["num"]==2:
        brk("Lunch Break","1 hour")
    elif t["num"]==1:
        brk("Tea Break","15 minutes")
    elif t["num"]==3:
        brk("Tea Break","15 minutes")

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("LO1 · Understood","The key concepts, principles and business value of design thinking."),
 ("LO2 · Applied","Design thinking methods to generate new ideas for your organisation."),
 ("LO3 · Uncovered","Real opportunities to apply design thinking where you work."),
 ("LO4 · Implemented","A plan to embed the five action phases across the organisation."),
 ("LO5 · Executed","A design concept as a tested low-fidelity prototype."),
 ("LO6 · Measured","The outcome of a design idea using metrics and a value/novelty map.")],
 kicker="LEARNING OUTCOMES",cols=2,size=14)
tile_grid("Take It Back to Work",[
 ("Pick one opportunity","Start with the impact-versus-effort winner from Activity 3."),
 ("Run a one-week sprint","Empathize and Define on days 1–2, Ideate day 3, Prototype and Test days 4–5."),
 ("Give it an owner","Use the RACI matrix so the concept has an accountable name against it."),
 ("Define the metric first","Agree what success looks like before you build anything."),
 ("Keep the ed-tools","All the canvases stay free and available in your browser after the course."),
 ("Share the result","Show colleagues the prototype and the user feedback — evidence beats opinion.")],
 kicker="NEXT STEPS",cols=2,size=13,accent=TEAL)
edtool_slide("Keep Using These Tools After the Course",
             "YOURS TO KEEP  ·  FREE  ·  NO LOGIN REQUIRED",
             descs=EDTOOL_TAKEAWAY)
content("Summary & Q&A",[
 "Design thinking is a human-centred, iterative approach to solving real problems.",
 "The five action phases — Empathize, Define, Ideate, Prototype, Test — loop rather than run once.",
 "Methods and visual tools make each phase concrete and repeatable.",
 "Metrics turn a good idea into a defensible business case.",
 "Questions?"],kicker="WRAP-UP")
assessment_slide()
flow_h("Assessment Flow",[
 "Scan the QR code on the LMS",
 "Scan the SSG QR code",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY",
 labels=["TRAQOM","Attendance","Assessment","Submit","Sign off"])
attendance_slide()
tile_grid("Recommended Courses",[
 "WSQ – Agile Project Management with Scrum",
 "WSQ – Digital Transformation Strategy for Business",
 "WSQ – Business Innovation with Emerging Technology",
 "WSQ – Data Visualisation for Business Decision Making",
 "WSQ – Customer Experience and Service Design",
 "WSQ – Business Continuity Management Essentials"],
 kicker="CONTINUE YOUR LEARNING",cols=2,size=15,accent=VIOLET)
tile_grid("Support",[
 ("✉  Email","enquiry@tertiaryinfotech.com"),
 ("☎  Telephone","+65 6100 0613"),
 ("🌐  Website","www.tertiarycourses.com.sg"),
 ("💬  After the course","Free post-course consultation on the subject matter is available to every learner.")],
 kicker="WE'RE HERE TO HELP  ·  DURING AND AFTER THE CLASS",cols=2,size=15,accent=BLUE)
big_statement("Thank You!","You can now apply design thinking to uncover real user needs, generate ideas, prototype them and prove they work.","KEEP DESIGNING",color=TEAL)

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
